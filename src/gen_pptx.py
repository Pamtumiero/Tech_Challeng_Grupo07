# -*- coding: utf-8 -*-
"""Deck executivo apresentacao_executiva.pptx - FOCADO em Geografia e Categorias."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figs"
OUT = ROOT / "apresentacao" / "apresentacao_executiva.pptx"

ROSA=RGBColor(0xDE,0x00,0x78); PRETO=RGBColor(0x1F,0x1F,0x1F); TEAL=RGBColor(0x00,0xA3,0xA1)
AMARELO=RGBColor(0xFF,0xB8,0x00); AZUL=RGBColor(0x5B,0x6B,0xF5); CINZA=RGBColor(0x8C,0x8C,0x8C)
BRANCO=RGBColor(0xFF,0xFF,0xFF); CLARO=RGBColor(0xF6,0xF3,0xF4)

prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height


def fill(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def box(s,l,t,w,h,text,size=18,color=PRETO,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,ls=1.05):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=ls
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold; r.font.name="Calibri"; r.font.color.rgb=color
    return tb
def rect(s,l,t,w,h,c):
    sp=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background(); return sp
def add_img(s,path,l,t,w,h):
    p=FIG/path
    try: iw,ih=Image.open(p).size
    except: s.shapes.add_picture(str(p),Inches(l),Inches(t),Inches(w),Inches(h)); return
    br=w/h; ir=iw/ih
    if ir>br: nw,nh=w,w/ir
    else: nh,nw=h,h*ir
    s.shapes.add_picture(str(p),Inches(l+(w-nw)/2),Inches(t+(h-nh)/2),Inches(nw),Inches(nh))

# ---- 1. Capa ----
s=prs.slides.add_slide(BLANK); fill(s,PRETO); rect(s,0,0,0.35,7.5,ROSA)
box(s,0.9,1.9,11.6,1.7,"Geografia e Categorias:\ndestravar o Brasil além do Sudeste",size=40,color=BRANCO,bold=True,ls=1.1)
box(s,0.95,4.0,11.5,0.6,"Análise de dados · Tech Challenge Fase 1 (POSTECH DTAT)",size=20,color=ROSA,bold=True)
box(s,0.95,4.8,11.5,0.5,"Brazilian E-Commerce Public Dataset by Olist · 64,6% da receita está no Sudeste",size=15,color=CINZA)
box(s,0.95,6.05,11.5,1.1,
    "Lucas Aguiar de Moura (RM375551) · Pamela Regina Tumiero da Costa (RM374909) · "
    "Guilherme Augusto Justino da Silva (RM374370)\nMariana Nogueira Salgado Zeron (RM374774) · Vitor Martino (RM374436)",
    size=12,color=BRANCO,ls=1.2)

# ---- 2. O problema + recorte ----
s=prs.slides.add_slide(BLANK); fill(s,BRANCO); rect(s,0,0,13.333,0.18,ROSA)
box(s,0.7,0.4,12,0.7,"O problema: um e-commerce do Sudeste",size=28,color=PRETO,bold=True)
add_img(s,"geo_00_panorama.png",0.4,1.3,7.6,5.6)
box(s,8.2,1.4,4.8,0.5,"As regiões críticas:",size=16,color=PRETO,bold=True)
for i,(t,d,c) in enumerate([
    ("Norte — restrição","Barreira geográfica (rios, estradas). Fora do alcance da empresa.",AZUL),
    ("NE / CO — foco","Demanda premium reprimida. Onde a solução tem retorno.",ROSA),
    ("Sul — benchmark","A região fora do Sudeste que deu certo.",TEAL)]):
    y=2.0+i*1.5; rect(s,8.2,y,0.16,1.25,c)
    box(s,8.5,y,4.5,0.5,t,size=15,color=PRETO,bold=True)
    box(s,8.5,y+0.45,4.5,0.8,d,size=12,color=CINZA)

# ---- 2.5 Adendo: escopo ----
s=prs.slides.add_slide(BLANK); fill(s,BRANCO); rect(s,0,0,13.333,0.18,PRETO)
box(s,0.7,0.4,12,0.8,"Adendo — usamos apenas pedidos entregues",size=26,color=PRETO,bold=True)
add_img(s,"geo_adendo_status.png",0.4,1.4,8.6,4.9)
box(s,9.0,1.55,4.0,5.2,
    "Escopo: só 'delivered'\n(96.478).\n\n"
    "Os 2.963 não entregues (FORA):\n"
    "• 1.107 perdidos no transporte\n• 625 cancelados\n• 609 sem estoque\n• 622 travados no funil\n\n"
    "São falhas reais (nota 1,3–2,0):\nR$ 370 mil de produto em risco.\n\n"
    "→ aponta 2º problema operacional\n(estoque + fulfillment).",
    size=13,color=PRETO,ls=1.16)

# ---- 3. Q1 ----
s=prs.slides.add_slide(BLANK); fill(s,BRANCO); rect(s,0,0,13.333,0.18,ROSA)
box(s,0.7,0.4,12.2,0.8,"Q1 — Por que NE e CO compram menos que o Sul?",size=26,color=PRETO,bold=True)
add_img(s,"geo_q1_oferta_frete.png",0.4,1.35,8.4,4.7)
box(s,0.7,6.2,12,1.1,"Falta OFERTA LOCAL, não demanda. NE tem 56 sellers e CO 79 — vs. 668 no Sul. Quase tudo vem de longe → "
    "frete ~31% do preço no NE → só o pedido caro compensa → as compras pequenas (que dão volume ao Sul) somem.",
    size=15,color=ROSA,bold=True,ls=1.15)

# ---- 4. Categorias ----
s=prs.slides.add_slide(BLANK); fill(s,BRANCO); rect(s,0,0,13.333,0.18,AZUL)
box(s,0.7,0.4,12,0.7,"Categorias confirmam: o frete molda o que cada região compra",size=23,color=PRETO,bold=True)
add_img(s,"geo_02_categoria_overindex.png",0.4,1.3,7.8,5.7)
box(s,8.4,1.6,4.6,5.0,
    "O Nordeste FOGE do volumoso:\n• Bed Bath Table −45%\n• Housewares −39%\n(frete proibitivo no pesado)\n\n"
    "E sobre-indexa no pequeno e caro:\n• Health Beauty +41%\n• Watches Gifts +22%\n(leve, caro — 'vale' o frete)\n\n"
    "→ a demanda premium reprimida\naparece no próprio sortimento.",
    size=15,color=PRETO,ls=1.18)

# ---- 5. Q2 ----
s=prs.slides.add_slide(BLANK); fill(s,BRANCO); rect(s,0,0,13.333,0.18,ROSA)
box(s,0.7,0.4,12.2,0.8,"Q2 — O gargalo da entrega é o transporte, não o seller",size=25,color=PRETO,bold=True)
add_img(s,"geo_q2_gargalo_entrega.png",0.4,1.35,8.0,5.0)
box(s,8.6,1.6,4.4,5.0,
    "O seller despacha rápido e por igual em todo o país (~2 dias).\n\n"
    "A diferença é 100% TRANSPORTE:\n6 dias no Sudeste →\n14 dias no Nordeste.\n\n"
    "NE tem o pior cumprimento de promessa (87%) — a perna longa é a mais variável.\n\n"
    "→ pressionar o seller não move o NE/CO.",
    size=15,color=PRETO,ls=1.18)

# ---- 5b. Q2 aprofundado: local + distancia ----
s=prs.slides.add_slide(BLANK); fill(s,BRANCO); rect(s,0,0,13.333,0.18,ROSA)
box(s,0.7,0.4,12.4,0.8,"A prova: vendedor local entrega — e o que pesa é a distância",size=23,color=PRETO,bold=True)
add_img(s,"geo_q2_distancia_ne.png",0.4,1.45,8.4,4.7)
box(s,8.9,1.55,4.1,5.2,
    "Local × não-local:\n• geral 8d vs 13d\n• NE 11d/91% vs 17d/87%\n• CO 7d/98% vs 13d/93%\n\n"
    "Por que o NE demora mesmo local?\nDISTÂNCIA — o NE é enorme.\n\n"
    "Mesmo estado → 7 dias\n(= Sudeste!)\nOutro estado do NE → 13 dias\n\n"
    "→ não um hub único:\nfulfillment por estado\n(BA, PE, CE).",
    size=14,color=PRETO,ls=1.16)

# ---- 6. Satisfacao ----
s=prs.slides.add_slide(BLANK); fill(s,BRANCO); rect(s,0,0,13.333,0.18,TEAL)
box(s,0.7,0.4,12,0.7,"A satisfação segue a promessa de entrega",size=26,color=PRETO,bold=True)
add_img(s,"geo_satisfacao_promessa.png",0.4,1.3,8.6,5.4)
box(s,9.1,1.6,3.9,5.2,
    "Baseline (linha de partida):\n\n🟥 Nordeste — 3,97\n(menor nota, maior alavanca)\n\n"
    "🟨 Centro-Oeste — 4,13\n🟩 Sul — 4,19 (= Sudeste)\n\n"
    "Welch t-test (no prazo × atraso):\nqueda de ~2 estrelas,\np ≈ 0 (significante).\n"
    "O atraso é o detrator nº 1.",
    size=14,color=PRETO,bold=False,ls=1.18)

# ---- 7. Sintese + recomendacoes ----
s=prs.slides.add_slide(BLANK); fill(s,BRANCO); rect(s,0,0,13.333,0.18,ROSA)
box(s,0.7,0.4,12,0.7,"Uma causa raiz, dois sintomas → o plano",size=26,color=PRETO,bold=True)
rect(s,0.7,1.35,12,1.15,CLARO)
box(s,0.9,1.5,11.6,0.9,"Q1 (compram menos) e Q2 (entrega atrasa) têm a MESMA causa: a oferta está longe do cliente. "
    "Colocar oferta no NE/CO baixa o frete (mais compras) E encurta o transporte (cumpre a promessa).",
    size=15,color=PRETO,bold=True,anchor=MSO_ANCHOR.MIDDLE,ls=1.15)
recs=[("1","Fulfillment distribuído por estado (BA, PE, CE)","não um hub único · Q1 frete + Q2 transporte · NE 3,97→4,15",ROSA),
      ("2","Recrutar sellers locais NE/CO","nas categorias que sobre-indexam · 56→150+ por estado",ROSA),
      ("3","Transportadora dedicada (corredor + intra-NE)","reduz variância do transporte · 90d",TEAL),
      ("4","Bundles regionais","categorias de sobre-indexação · sobe ticket/volume",TEAL),
      ("5","Recalibrar data prometida por rota","tático, até a oferta local crescer",TEAL)]
y=2.75
for n,t,d,c in recs:
    rect(s,0.7,y,0.5,0.72,c); box(s,0.7,y,0.5,0.72,n,size=20,color=BRANCO,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    box(s,1.4,y,7.2,0.45,t,size=16,color=PRETO,bold=True)
    box(s,1.4,y+0.4,11.4,0.35,d,size=12,color=CINZA)
    y+=0.86

# ---- 7b. Simulacao 30% ----
s=prs.slides.add_slide(BLANK); fill(s,BRANCO); rect(s,0,0,13.333,0.18,TEAL)
box(s,0.7,0.4,12.2,0.8,"Quanto isso vale: simulação do cenário 30%",size=26,color=PRETO,bold=True)
box(s,0.7,1.35,12,1.0,"Cenário-base: migrar 30% dos pedidos hoje despachados de FORA da região para sourcing local "
    "(aplicando o desempenho já observado nos pedidos locais). Números próprios — cenário, não previsão.",
    size=15,color=PRETO,ls=1.18)
sc=[("R$ 51 mil","economia de frete/ano","NE R$ 35,8k + CO R$ 15,4k",ROSA),
    ("169","atrasos evitados","NE 92 + CO 77",TEAL),
    ("4.209","pedidos otimizados","migração de 30% dos fluxos",PRETO),
    ("+0,3★","nota dos migrados","de 'não-local' p/ 'local'",ROSA)]
x=0.7
for v,l,nt,c in sc:
    rect(s,x,2.7,2.9,1.7,CLARO)
    box(s,x,2.95,2.9,0.7,v,size=26,color=c,bold=True,align=PP_ALIGN.CENTER)
    box(s,x+0.1,3.7,2.7,0.4,l,size=13,color=PRETO,bold=True,align=PP_ALIGN.CENTER)
    box(s,x+0.1,4.05,2.7,0.3,nt,size=10,color=CINZA,align=PP_ALIGN.CENTER)
    x+=3.05
box(s,0.7,4.9,12,1.2,"A economia de frete financia a própria atração de sellers locais (recomendação #2). "
    "E os pedidos migrados ganham satisfação — que, sem recompra para diluir o CAC, é o que protege a primeira compra.",
    size=15,color=ROSA,bold=True,ls=1.2)

# ---- 8. Fechamento ----
s=prs.slides.add_slide(BLANK); fill(s,PRETO); rect(s,0,0,0.35,7.5,ROSA)
box(s,0.9,1.7,11.6,1.5,"O mapa de receita não é destino —\né consequência de onde a oferta está.",size=30,color=BRANCO,bold=True,ls=1.1)
box(s,0.95,3.9,11.5,1.6,"NE e CO compram pouco E recebem atrasado pela mesma razão: não há vendedor por perto.\n\n"
    "Aproximar a oferta — distribuída nos estados de maior demanda, não num hub único — converte a demanda premium e cumpre a promessa — "
    "medido pela satisfação dessas regiões.",size=18,color=CLARO,ls=1.25)
box(s,0.95,6.4,11.5,0.5,"Meta-norte: satisfação média do NE 3,97 → 4,15 · Centro-Oeste 4,13 → 4,20",size=15,color=ROSA,bold=True)

prs.save(str(OUT))
print(f"OK -> {OUT.relative_to(ROOT)}  ({len(prs.slides._sldIdLst)} slides, foco geográfico)")
